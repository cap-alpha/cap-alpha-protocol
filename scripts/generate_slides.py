from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_deck():
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # --- Slide 1: Title Slide ---
    slide1 = prs.slides.add_slide(title_slide_layout)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = "Cap Alpha Protocol: The Future of Roster Capital"
    subtitle1.text = "A Quantitative Intelligence Engine for the $200B NFL Market\n\nExecutive Briefing"

    # --- Slide 2: The Inefficiency (The Conflict) ---
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes2 = slide2.shapes
    title2 = shapes2.title
    body2 = shapes2.placeholders[1]
    title2.text = "The Inefficiency (The Conflict)"
    tf2 = body2.text_frame
    
    p = tf2.add_paragraph()
    p.text = "The Problem: NFL GMs manage $250M+ annual budgets using anecdotal scouting and lagged media metrics."
    p.level = 0
    p = tf2.add_paragraph()
    p.text = "The Liability: A single $50M error in dead cap destroys a franchise's championship window."
    p.level = 0
    p = tf2.add_paragraph()
    p.text = "The Solution: Institutional-grade, real-time capital velocity tracking to mitigate catastrophic financial risk."
    p.level = 0

    # --- Slide 3: The Cap Alpha Oracle (Real-Time Hydration) ---
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    shapes3 = slide3.shapes
    title3 = shapes3.title
    body3 = shapes3.placeholders[1]
    title3.text = "The Cap Alpha Oracle (Real-Time Hydration)"
    tf3 = body3.text_frame
    
    p = tf3.add_paragraph()
    p.text = "Evolution from 'Hobbyist Scripts' into a 'Regulatory-Grade Data Pipeline.'"
    p.level = 0
    p = tf3.add_paragraph()
    p.text = "A fully autonomous, serverless engine driving unparalleled execution speed."
    p.level = 0
    p = tf3.add_paragraph()
    p.text = "Scans the Top 500 NFL Assets daily, monitoring for off-field volatility, hidden injury risk, and contract anomalies."
    p.level = 0

    # --- Slide 4: Regulatory-Grade Intelligence (Signal over Noise) ---
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    shapes4 = slide4.shapes
    title4 = shapes4.title
    body4 = shapes4.placeholders[1]
    title4.text = "Regulatory-Grade Intelligence"
    tf4 = body4.text_frame
    
    p = tf4.add_paragraph()
    p.text = "Unstructured internet chatter is a liability; accountable intelligence is a proprietary asset."
    p.level = 0
    p = tf4.add_paragraph()
    p.text = "Immutable Auditability: Predictions and signals are hashed to a verifiable cryptographic ledger."
    p.level = 0
    p = tf4.add_paragraph()
    p.text = "We cannot hide our misses. This guarantees 100% honesty in our historical FMV projections."
    p.level = 0
    p = tf4.add_paragraph()
    p.text = "Institutional Confidence: Verifiable audit trails and source attribution for Front Office review."
    p.level = 0

    # --- Slide 5: The Infrastructure Strategy ---
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    shapes5 = slide5.shapes
    title5 = shapes5.title
    body5 = shapes5.placeholders[1]
    title5.text = "The Infrastructure Strategy (Efficiency & Execution)"
    tf5 = body5.text_frame
    
    p = tf5.add_paragraph()
    p.text = "MotherDuck Integration: Deployed a sub-second latency cloud data warehouse."
    p.level = 0
    p = tf5.add_paragraph()
    p.text = "Lean Operations: Minimizing burn rate while maximizing execution speed through a decoupled serverless architecture."
    p.level = 0
    p = tf5.add_paragraph()
    p.text = "B2B Ready: Enterprise-grade authentication (Clerk) explicitly gating proprietary Roster Intel."
    p.level = 0

    # --- Slide 6: Capturing Value (Monetization & Velocity) ---
    slide6 = prs.slides.add_slide(bullet_slide_layout)
    shapes6 = slide6.shapes
    title6 = shapes6.title
    body6 = shapes6.placeholders[1]
    title6.text = "Capturing Value (Monetization & Velocity)"
    tf6 = body6.text_frame
    
    p = tf6.add_paragraph()
    p.text = "Enterprise Play (B2B): Executive dashboards actively marketed as 'Insurance' for General Managers."
    p.level = 0
    p = tf6.add_paragraph()
    p.text = "Consumer Play (B2C): Premium subscription-access for the sophisticated, data-driven NFL bettor."
    p.level = 0
    p = tf6.add_paragraph()
    p.text = "The Next Move: Igniting the autonomous Reinforcement Learning flywheel to scale predictions automatically."
    p.level = 0

    prs.save('cap_alpha_sprint_12_showcase.pptx')
    print("Presentation saved as cap_alpha_sprint_12_showcase.pptx")

if __name__ == '__main__':
    create_deck()
